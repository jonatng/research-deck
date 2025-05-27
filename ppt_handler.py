from pptx import Presentation
from pptx.util import Inches

def create_ppt_with_summaries(summaries, output_file="summaries.pptx"):
    prs = Presentation()
    for url, summary in summaries.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = url
        slide.placeholders[1].text = summary
    prs.save(output_file)

def append_to_existing_ppt(pptx_path, summaries, output_file="appended.pptx"):
    prs = Presentation(pptx_path)
    for url, summary in summaries.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = url
        slide.placeholders[1].text = summary
    prs.save(output_file)
