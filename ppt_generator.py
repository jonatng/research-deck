from pptx import Presentation
from pptx.util import Inches
import os

def create_presentation():
    return Presentation()

def add_summary_slide(prs, summary, title="Summary"):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = summary
    return prs

def save_presentation(prs, filename="summary_output.pptx"):
    prs.save(filename)
    return filename

def append_to_existing_ppt(summary, ppt_path):
    prs = Presentation(ppt_path)
    add_summary_slide(prs, summary, title="Appended Summary")
    new_path = ppt_path.replace(".pptx", "_updated.pptx")
    prs.save(new_path)
    return new_path
